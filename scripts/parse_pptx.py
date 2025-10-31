import os, json, re, sys
from glob import glob
from tqdm import tqdm
from pptx import Presentation

RAW_DIR = os.path.join("data_raw", "aws_lectures")
OUT_FILE = os.path.join("data_processed", "chunks.jsonl")

def clean(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\u2022", "- ").replace("\u25cf", "- ")
    text = re.sub(r"\r\n|\r", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

def extract_shape_text(shape) -> str:
    # text frames
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        parts = []
        for para in shape.text_frame.paragraphs:
            runs = [run.text for run in para.runs if run.text]
            parts.append("".join(runs) if runs else para.text)
        return "\n".join([p for p in parts if p])

    # simple fallback
    if hasattr(shape, "text") and shape.text:
        return str(shape.text)

    # tables (shape_type 19)
    if getattr(shape, "shape_type", None) == 19 and hasattr(shape, "table"):
        rows = []
        tbl = shape.table
        for r in tbl.rows:
            row_txt = []
            for c in r.cells:
                row_txt.append(clean(c.text))
            rows.append(" | ".join(row_txt))
        if rows:
            header = rows[0]
            sep = " | ".join(["---"] * len(rows[0].split(" | ")))
            body = "\n".join(rows[1:]) if len(rows) > 1 else ""
            return f"{header}\n{sep}\n{body}".strip()

    return ""

def extract_slide_text(slide) -> str:
    body_bits = []
    for shp in slide.shapes:
        try:
            t = extract_shape_text(shp)
            if t:
                body_bits.append(t)
        except Exception:
            continue
    body = "\n".join([b for b in body_bits if b])

    notes_txt = ""
    if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
        notes_parts = []
        for p in slide.notes_slide.notes_text_frame.paragraphs:
            runs = [run.text for run in p.runs if run.text]
            notes_parts.append("".join(runs) if runs else p.text)
        notes_txt = "\n".join([p for p in notes_parts if p])

    combined = []
    if body.strip():
        combined.append(body)
    if notes_txt.strip():
        combined.append(f"[Speaker Notes]\n{notes_txt}")
    return clean("\n\n".join(combined))

def parse_pptx_file(path: str):
    prs = Presentation(path)
    slides_chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        text = extract_slide_text(slide)
        slides_chunks.append({
            "source": os.path.basename(path),
            "source_path": path,
            "doc_type": "pptx",
            "slide": i,
            "text": text
        })
    return slides_chunks

def main():
    pptx_files = sorted(glob(os.path.join(RAW_DIR, "*.pptx")))
    if not pptx_files:
        print(f"No PPTX files found in {RAW_DIR}.", file=sys.stderr)
        sys.exit(1)

    os.makedirs(os.path.dirname(OUT_FILE), exist_ok=True)

    total = 0
    with open(OUT_FILE, "w", encoding="utf-8") as out:
        for f in tqdm(pptx_files, desc="Parsing PPTX"):
            try:
                chunks = parse_pptx_file(f)
                for c in chunks:
                    c["text"] = c.get("text", "")
                    out.write(json.dumps(c, ensure_ascii=False) + "\n")
                total += len(chunks)
            except Exception as e:
                print(f"[WARN] Failed to parse {f}: {e}", file=sys.stderr)

    print(f"Done. Wrote {total} slide chunks to {OUT_FILE}")

if __name__ == "__main__":
    main()

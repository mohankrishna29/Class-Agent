# scripts/ingest.py
from __future__ import annotations
import os, re, uuid, hashlib, json
from pathlib import Path
from typing import List, Dict, Any
from dataclasses import dataclass

import backoff
import numpy as np
import pandas as pd
from tqdm import tqdm
from dotenv import load_dotenv

# Extractors
import fitz               # PyMuPDF
import pptx               # python-pptx
import docx               # python-docx
from PIL import Image
import pytesseract

# Vectors
import faiss
from openai import OpenAI

load_dotenv()
ROOT = Path(__file__).resolve().parents[1]

# --- Paths (keep your existing layout) ---
RAW_ROOTS = [
    ROOT / "data_raw" / "course_assets",           # includes canvas/<run-id>/ from canvas_extract.py
    # You can add more roots here if needed
]
CLEAN = ROOT / "data_clean"
INDEX_DIR = ROOT / "index"
MANIFESTS = ROOT / "manifests"

CLEAN.mkdir(parents=True, exist_ok=True)
INDEX_DIR.mkdir(parents=True, exist_ok=True)
MANIFESTS.mkdir(parents=True, exist_ok=True)

# --- Config ---
EMBED_PROVIDER = os.getenv("EMBED_PROVIDER", "openai")  # "openai" or "local"
EMBED_LOCAL_MODEL = os.getenv("EMBED_LOCAL_MODEL", "thenlper/gte-small")  # good 384-d model
_local_model = None
SKIP_OCR = os.getenv("SKIP_OCR", "true").lower() == "true"
EMBED_MODEL = os.getenv("EMBED_MODEL", "text-embedding-3-large")
BATCH = int(os.getenv("EMBED_BATCH", "512"))
LEAK_ANSWERS = os.getenv("LEAK_ANSWERS", "false").lower() == "true"  # ignored here unless you bring question banks back

SUPPORTED = {".pdf", ".pptx", ".docx", ".md", ".txt", ".png", ".jpg", ".jpeg", ".html", ".htm"}

client = OpenAI()


# --------- Utilities ----------
def sha256_bytes(b: bytes) -> str:
    return hashlib.sha256(b).hexdigest()

def sha256_text(s: str) -> str:
    return sha256_bytes(s.encode("utf-8"))

def write_jsonl(path: Path, objs: List[dict]):
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        for o in objs:
            f.write(json.dumps(o, ensure_ascii=False) + "\n")

def approx_chunk(text: str, N=3500, O=800):
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    i = 0
    while i < len(text):
        piece = text[i:i+N]
        if not piece.strip():
            break
        yield piece
        i += (N - O)


# --------- Embeddings ----------
@backoff.on_exception(backoff.expo, Exception, max_time=120)
def embed_batch(texts: List[str]) -> List[List[float]]:
    global _local_model
    if EMBED_PROVIDER == "openai":
        resp = client.embeddings.create(model=EMBED_MODEL, input=texts)
        return [d.embedding for d in resp.data]
    else:
        # local
        if _local_model is None:
            from sentence_transformers import SentenceTransformer
            _local_model = SentenceTransformer(EMBED_LOCAL_MODEL)  # 384 dims
        # returns numpy array [n, dim]
        X = _local_model.encode(texts, normalize_embeddings=True, batch_size=64, show_progress_bar=False)
        return [x.tolist() for x in X]



# --------- Extractors ----------
def extract_pdf(path: Path) -> List[dict]:
    doc = fitz.open(path)
    out = []
    for i, page in enumerate(doc, start=1):
        txt = page.get_text("text")
        if not txt.strip():
            # OCR fallback
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            txt = pytesseract.image_to_string(img)
        if not txt.strip():
            continue
        out.append(_mk_chunk(
            text=txt.strip(),
            path=path,
            locator=f"page={i}",
            mimetype="application/pdf",
            extra={"page": i},
        ))
    return out

def extract_pptx(path: Path) -> List[dict]:
    prs = pptx.Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            texts.append(slide.notes_slide.notes_text_frame.text or "")
        txt = "\n".join(t.strip() for t in texts if t and t.strip())
        if not txt:
            continue
        out.append(_mk_chunk(
            text=txt,
            path=path,
            locator=f"slide={i}",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            extra={"slide": i},
        ))
    return out

def extract_docx(path: Path) -> List[dict]:
    d = docx.Document(path)
    full = [p.text for p in d.paragraphs if p.text and p.text.strip()]
    text = "\n".join(full)
    return _chunk_textlike(text, path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")

def extract_md_txt(path: Path) -> List[dict]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    mime = "text/markdown" if path.suffix.lower() == ".md" else "text/plain"
    return _chunk_textlike(text, path, mime)

def extract_html(path: Path) -> List[dict]:
    # If you ran canvas_extract.py, many pages already have .txt siblings.
    # Still, we provide a light HTML → text fallback.
    raw = path.read_text(encoding="utf-8", errors="ignore")
    # crude strip; for full fidelity rely on the .txt created by canvas_extract.py
    cleaned = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", raw)
    cleaned = re.sub(r"(?is)<[^>]+>", "\n", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return _chunk_textlike(cleaned, path, "text/html")

def extract_image(path: Path) -> List[dict]:
    img = Image.open(path)
    txt = pytesseract.image_to_string(img)
    if not txt.strip():
        return []
    return [_mk_chunk(
        text=txt.strip(),
        path=path,
        locator="image=1",
        mimetype=f"image/{path.suffix.lower().lstrip('.')}",
    )]

def _chunk_textlike(text: str, path: Path, mimetype: str) -> List[dict]:
    out = []
    for off, piece in enumerate(approx_chunk(text), start=0):
        out.append(_mk_chunk(
            text=piece,
            path=path,
            locator=f"offset={off}",
            mimetype=mimetype,
        ))
    return out

def _mk_chunk(*, text: str, path: Path, locator: str, mimetype: str, extra: Dict[str, Any] | None = None) -> dict:
    rel = path.relative_to(ROOT)
    return {
        "chunk_id": str(uuid.uuid4()),
        "text": text,
        "source_id": str(rel),
        "locator": locator,
        "mimetype": mimetype,
        "title": path.stem,
        "extra": extra or {},
        "hash": sha256_text(text),
    }

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".docx": extract_docx,
    ".md": extract_md_txt,
    ".txt": extract_md_txt,
    ".html": extract_html,
    ".htm": extract_html,
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
}


# --------- Manifest ----------
def current_files() -> pd.DataFrame:
    rows = []
    for root in RAW_ROOTS:
        if not root.exists():
            continue
        for p in root.rglob("*"):
            if p.is_file() and p.suffix.lower() in SUPPORTED:
                rows.append({
                    "path": str(p.relative_to(ROOT)),
                    "sha256": sha256_bytes(p.read_bytes()),
                    "mtime": p.stat().st_mtime,
                })
    df = pd.DataFrame(rows).sort_values("path").reset_index(drop=True)
    return df

def load_manifest() -> pd.DataFrame | None:
    f = MANIFESTS / "files.jsonl"
    if f.exists():
        return pd.read_json(f, lines=True)
    return None

def save_manifest(df: pd.DataFrame):
    out = MANIFESTS / "files.jsonl"
    out.unlink(missing_ok=True)
    df.to_json(out, orient="records", lines=True)


# --------- Ingest pipeline ----------
def build_clean_if_needed(df: pd.DataFrame):
    """
    For each RAW file, if there's no corresponding CLEAN jsonl,
    run the right extractor and write one jsonl with all chunks.
    """
    for _, row in tqdm(df.iterrows(), total=len(df), desc="Extracting"):
        rel = Path(row["path"])
        src = ROOT / rel
        clean_out = CLEAN / (rel.with_suffix(".jsonl").name)

        if clean_out.exists():
            # already extracted → skip
            continue

        ext = src.suffix.lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            continue

        try:
            chunks = extractor(src)
            # Light garbage filter
            filtered = []
            for ch in chunks:
                t = ch["text"].strip()
                if len(t) < 120:
                    continue
                if len(set(t)) < 10:
                    continue
                filtered.append(ch)

            if not filtered:
                # still write an empty file to mark "processed", prevents re-work
                write_jsonl(clean_out, [])
            else:
                write_jsonl(clean_out, filtered)

        except Exception as e:
            print(f"[WARN] Failed to extract {src}: {e}")
            # do not create the clean file; next run will retry


def load_all_clean() -> tuple[list[str], list[str], list[dict]]:
    """
    Load (ids, texts, metas) from:
      - data_clean/**                (Canvas chunks)
      - data_processed/**            (previous run, any depth: includes chunks*.jsonl)
    Accepts JSONL or JSON array. Tolerates LangChain-style schema
    (page_content + metadata) and synthesizes chunk_id when missing.
    """
    import json, hashlib
    ids, texts, metas = [], [], []

    SEARCH_ROOTS = [ROOT / "data_clean", ROOT / "data_processed"]

    def take_text(obj: dict) -> str | None:
        # common variants
        t = (obj.get("text")
             or obj.get("page_content")
             or obj.get("content")
             or obj.get("body"))
        if t is None and isinstance(obj.get("data"), dict):
            t = obj["data"].get("text")
        return (t or "").strip()

    def ensure_chunk_id(obj: dict) -> str:
        cid = obj.get("chunk_id") or obj.get("id") or obj.get("uuid")
        if not cid:
            cid = hashlib.sha256((obj.get("text") or "").encode("utf-8")).hexdigest()
            obj["chunk_id"] = cid
        return cid

    def normalize_meta(obj: dict, src_path: Path) -> dict:
        # If there's a nested metadata block (LangChain), merge it
        md = obj.get("metadata") or {}
        if isinstance(md, dict):
            for k, v in md.items():
                obj.setdefault(k, v)

        # Fill standard fields
        obj.setdefault("source_id",
                       obj.get("source")
                       or obj.get("file")
                       or obj.get("document_id")
                       or str(src_path))
        obj.setdefault("locator",
                       obj.get("locator")
                       or obj.get("loc")
                       or (f"page={obj['page']}" if obj.get("page") is not None else "offset=0"))
        obj.setdefault("mimetype", obj.get("mimetype") or obj.get("mime") or "text/plain")
        obj.setdefault("title", obj.get("title") or Path(str(obj.get("source_id"))).stem)
        obj.setdefault("extra", obj.get("extra") or {})
        return obj

    def consume_obj(obj: dict, jf: Path):
        t = take_text(obj)
        if not t:
            return
        obj["text"] = t  # normalize
        obj = normalize_meta(obj, jf)
        cid = ensure_chunk_id(obj)
        ids.append(cid)
        texts.append(t)
        metas.append(obj)

    for base in SEARCH_ROOTS:
        if not base.exists():
            continue
        files = list(base.rglob("*.jsonl")) + list(base.rglob("*.json"))
        for jf in sorted(files):
            try:
                with jf.open("r", encoding="utf-8") as f:
                    first = f.read(1); f.seek(0)
                    if first == "[":
                        arr = json.load(f)
                        if isinstance(arr, list):
                            for obj in arr:
                                if isinstance(obj, dict): consume_obj(obj, jf)
                    else:
                        for line in f:
                            line = line.strip()
                            if not line: continue
                            obj = json.loads(line)
                            if isinstance(obj, dict): consume_obj(obj, jf)
            except Exception as e:
                print(f"[WARN] Skipping unreadable file {jf}: {e}")

    return ids, texts, metas



def build_index():
    """
    Embed all cleaned chunks (from data_clean/ + data_processed/),
    resume safely via shard files, then build a FAISS IP (cosine) index.
    """
    ids, texts, metas = load_all_clean()
    if not texts:
        print("No cleaned chunks found. Did extraction run?")
        return

    shard_dir = INDEX_DIR / "shards"
    shard_dir.mkdir(parents=True, exist_ok=True)

    def shard_path(idx: int) -> Path:
        return shard_dir / f"part-{idx:05d}.npy"

    n = len(texts)
    num_batches = (n + BATCH - 1) // BATCH

    # ---- Embed in shards (resume-friendly) ----
    for bi in tqdm(range(num_batches), desc="Embedding"):
        sp = shard_path(bi)
        if sp.exists():
            continue  # already computed
        start = bi * BATCH
        end = min(start + BATCH, n)
        batch = texts[start:end]
        emb = embed_batch(batch)                 # returns list[list[float]]
        arr = np.asarray(emb, dtype="float32")   # shape [m, dim]
        np.save(sp, arr)

    # ---- Load shards and stack ----
    parts = []
    for bi in range(num_batches):
        sp = shard_path(bi)
        if not sp.exists():
            raise RuntimeError(f"Missing shard {sp}. Re-run ingest to complete embeddings.")
        parts.append(np.load(sp))
    X = np.vstack(parts).astype("float32")
    assert X.shape[0] == n, f"Vectors {X.shape[0]} != texts {n}"

    # Normalize for cosine similarity via inner product
    faiss.normalize_L2(X)

    # Build FAISS index with correct dimension
    dim = X.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(X)

    # Write artifacts
    faiss.write_index(index, str(INDEX_DIR / "faiss.index"))
    pd.DataFrame(metas).to_parquet(INDEX_DIR / "meta.parquet", index=False)

    print(f"Indexed {len(metas)} chunks (dim={dim}) → {INDEX_DIR/'faiss.index'}")
    print(f"Meta → {INDEX_DIR/'meta.parquet'}")



def main():
    # 1) Scan raw roots → manifest
    df = current_files()
    save_manifest(df)

    # 2) Create clean jsonl per source (skip if already exists)
    build_clean_if_needed(df)

    # 3) Build/overwrite the FAISS index from all clean chunks
    build_index()


if __name__ == "__main__":
    main()
